from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document

def extract_text(file_path):
    if file_path.endswith('.pdf'):
        return extract_pdf(file_path)
    elif file_path.endswith('.pptx'):
        return extract_ppt(file_path)
    elif file_path.endswith('.docx'):
        return extract_docx(file_path)
    elif file_path.endswith('.txt'):
        return extract_txt(file_path)
    return ""

def extract_pdf(path):
    text = ""
    try:
        with open(path, 'rb') as file:
            pdf_reader = PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + " "
    except Exception as e:
        print(f"Error extracting PDF: {e}")
    return text

def extract_ppt(path):
    text = ""
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
    except Exception as e:
        print(f"Error extracting PPT: {e}")
    return text

def extract_docx(path):
    text = ""
    try:
        doc = Document(path)
        for para in doc.paragraphs:
            text += para.text + " "
    except Exception as e:
        print(f"Error extracting DOCX: {e}")
    return text

def extract_txt(path):
    text = ""
    try:
        with open(path, 'r', encoding='utf-8') as file:
            text = file.read()
    except Exception as e:
        print(f"Error extracting TXT: {e}")
    return text
